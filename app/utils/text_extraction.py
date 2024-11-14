from io import BytesIO
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import re
import openpyxl
from langchain_community.document_loaders import WebBaseLoader
import base64
from openai import OpenAI
from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled
from app.config import config
client = OpenAI(api_key=config.OPENAI_API_KEY)

def get_pdf_text(pdf_doc):
    text = ""
    pdf_data = BytesIO(pdf_doc)
    pdf_reader = PdfReader(pdf_data)
    for page in pdf_reader.pages:
        text += page.extract_text()
    cleaned_text = re.sub(r'\s+', ' ', text)
    cleaned_text = cleaned_text.strip()
    return cleaned_text
    return text


def get_docx_text(docx_file):
    text = ""
    docx_data = BytesIO(docx_file)
    document = Document(docx_data)
    # Extract text from each paragraph
    for paragraph in document.paragraphs:
        text += paragraph.text + " "
    # Clean up extra whitespace and newlines
    cleaned_text = re.sub(r'\s+', ' ', text).strip()
    return cleaned_text

def get_pptx_text(pptx_file):
    text = ""
    pptx_data = BytesIO(pptx_file)
    presentation = Presentation(pptx_data)
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + " "
            elif shape.has_table:  # Check for tables
                for row in shape.table.rows:
                    for cell in row.cells:
                        text += cell.text + " "

    # Clean up extra whitespace and newlines
    cleaned_text = re.sub(r'\s+', ' ', text).strip()
    return cleaned_text

def get_xlsx_text(xlsx_file):
    text = ""
    # Load the Excel workbook
    xlsx_data = BytesIO(xlsx_file)
    workbook = openpyxl.load_workbook(xlsx_data, data_only=True)
    # Loop through all sheets in the workbook
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        # Loop through all rows and columns to extract text
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.value:
                    text += str(cell.value) + " "
    # Clean up extra whitespace and newlines
    cleaned_text = re.sub(r'\s+', ' ', text).strip()
    return cleaned_text

def get_web_contents(url):
    loader = WebBaseLoader(f"{url}")
    scraped_contents = loader.load()
    for document in scraped_contents:
        page_content = document.page_content
        
        # Clean the page_content
        cleaned_content = page_content.strip()
        cleaned_content = ' '.join(cleaned_content.split())
        
    return cleaned_content

def get_image_text(image_data: bytes):
    # Encode image in base64 using the raw byte data
    base64_image = base64.b64encode(image_data).decode('utf-8')
    
    # Send request to OpenAI for a detailed description
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Write a detailed description of this image, including subtle details."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                ],
            }
        ],
    )

    # Extract and return the description
    return response.choices[0].message.content

def extract_yt_transcript(video_url):
    """Extracts the transcript from a YouTube video."""
    try:
        # Extract the video ID from the URL
        video_id = video_url.split("=")[1]
        
        # Try to get the transcript list for the video
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        
        # Check if there are any available transcripts
        for transcript in transcript_list:
            transcript_text_list = transcript.fetch()
            lang = transcript.language
            transcript_text = ""
            
            # If English transcript is available, return it
            if transcript.language_code == 'en':
                for line in transcript_text_list:
                    transcript_text += " " + line["text"]
                return transcript_text
            
            # If a translatable transcript is available, return the translated English transcript
            elif transcript.is_translatable:
                english_transcript_list = transcript.translate('en').fetch()
                for line in english_transcript_list:
                    transcript_text += " " + line["text"]
                return transcript_text
        
        # If no English transcript or translatable transcript is found, raise an error
        raise ValueError("No transcript available for this video.")
    
    except ValueError as ve:
        raise ve  # Raise the custom error if no transcript found
    except TranscriptsDisabled:
        raise ValueError("Transcript is disabled for this video.")
    except Exception as e:
        raise ValueError(f"An error occurred: {str(e)}")

        